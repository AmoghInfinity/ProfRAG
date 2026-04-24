import os
import re
from typing import List, Dict
from pptx import Presentation


def clean_text(text: str) -> str:
    """
    Normalize and clean text.
    """
    if not text:
        return ""

    text = re.sub(r"\s+", " ", text)
    return text.strip()


def is_noise_line(line: str) -> bool:
    """
    Detect generic noise like headers/footers.
    Generalized (no hardcoding to specific files).
    """
    line = line.strip()

    if not line:
        return True

    # Pure numbers (page/slide numbers)
    if re.fullmatch(r"\d+", line):
        return True

    # High digit ratio (e.g., "Module-II 100")
    digits = sum(c.isdigit() for c in line)
    if len(line) > 0 and digits / len(line) > 0.3:
        return True

    # Repeated symbols
    if re.fullmatch(r"[-–—_=]+", line):
        return True

    # Academic header keywords (general, not specific)
    if re.search(r"(module|lecture|chapter|unit)", line, re.IGNORECASE):
        return True

    # Mostly uppercase short lines (common headers)
    words = line.split()
    if len(words) <= 5 and all(w.isupper() for w in words if w.isalpha()):
        return True

    # Very short lines
    if len(line) < 3:
        return True

    return False


def extract_text_from_ppt(file_path: str) -> List[Dict]:
    """
    Extract structured text from PPT in LangChain-friendly format.
    """

    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    presentation = Presentation(file_path)
    documents = []

    file_name = os.path.basename(file_path)

    for i, slide in enumerate(presentation.slides):
        title = ""
        content_parts = []

        # Extract title
        if slide.shapes.title:
            title = clean_text(slide.shapes.title.text)

        # Extract content
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                raw_text = shape.text
                lines = [clean_text(l) for l in raw_text.split("\n") if clean_text(l)]

                for line in lines:
                    if not line or line == title:
                        continue

                    if is_noise_line(line):
                        continue

                    content_parts.append(line)

        full_text = f"{title}. {' '.join(content_parts)}".strip()
        full_text = re.sub(r"\s+", " ", full_text)

        if not full_text:
            continue

        documents.append({
            "text": full_text,
            "metadata": {
                "source": "ppt",
                "file_name": file_name,
                "slide": i + 1,
                "title": title
            }
        })

    return documents


if __name__ == "__main__":
    sample_path = "data/raw/sample.pptx"
    docs = extract_text_from_ppt(sample_path)

    for doc in docs[:3]:
        print("\n--- Document ---")
        print(doc["text"][:200])
        print(doc["metadata"])